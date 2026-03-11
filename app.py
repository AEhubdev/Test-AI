import streamlit as st
import pandas as pd
import openai
import io
import re
import sys
import matplotlib.pyplot as plt
import seaborn as sns
from PyPDF2 import PdfReader
from pptx import Presentation

# --- 1. Page Config ---
st.set_page_config(page_title="Ultimate AI Data Scientist", layout="wide", page_icon="📊")

# --- 2. API Key Gatekeeper ---
if "openai_key" not in st.session_state:
    st.session_state.openai_key = None

if st.session_state.openai_key is None:
    st.title("🛡️ Enterprise AI Data Agent")
    user_key = st.text_input("Enter OpenAI API Key", type="password")
    if st.button("Unlock System"):
        if user_key.startswith("sk-"):
            st.session_state.openai_key = user_key
            st.rerun()
    st.stop()

client = openai.OpenAI(api_key=st.session_state.openai_key)
MODEL = "gpt-4o"

# --- 3. Session State Initialization ---
if "files" not in st.session_state:
    st.session_state.files = {}
if "history" not in st.session_state:
    st.session_state.history = []


# --- 4. File Parsers ---
def load_file(uploaded_file):
    name = uploaded_file.name
    try:
        if name.endswith(('.xlsx', '.xls')):
            return pd.read_excel(uploaded_file)
        elif name.endswith('.csv'):
            return pd.read_csv(uploaded_file)
        elif name.endswith('.pdf'):
            return "\n".join([p.extract_text() for p in PdfReader(uploaded_file).pages if p.extract_text()])
        elif name.endswith('.pptx'):
            return "\n".join(
                [s.text for sl in Presentation(uploaded_file).slides for s in sl.shapes if hasattr(s, "text")])
    except Exception as e:
        return f"Error: {e}"
    return None


# --- 5. Sidebar: File Uploads & Management ---
with st.sidebar:
    st.header("📂 File Center")
    uploads = st.file_uploader("Upload Data/Docs", type=["csv", "xlsx", "pdf", "pptx"], accept_multiple_files=True)
    if uploads:
        for f in uploads:
            if f.name not in st.session_state.files:
                st.session_state.files[f.name] = load_file(f)

    if st.button("Reset Everything", type="primary"):
        st.session_state.files = {}
        st.session_state.history = []
        st.rerun()

    st.divider()
    if st.session_state.files:
        st.subheader("📥 Export Changes")
        for fn, content in st.session_state.files.items():
            if isinstance(content, pd.DataFrame):
                buf = io.BytesIO()
                content.to_excel(buf, index=False)
                st.sidebar.download_button(f"Download {fn}", buf.getvalue(), f"edited_{fn}.xlsx",
                                           use_container_width=True)

# --- 6. Main Workspace: Real-time Previews ---
st.title("🤖 AI Multi-Modal Data Scientist")

if st.session_state.files:
    st.subheader("📋 Live Data Preview")
    # This creates a tab for every file so you can see exactly what is happening
    tabs = st.tabs([f"📄 {fn}" for fn in st.session_state.files.keys()])
    for i, (fn, content) in enumerate(st.session_state.files.items()):
        with tabs[i]:
            if isinstance(content, pd.DataFrame):
                st.dataframe(content, height=300, use_container_width=True)
                st.caption(f"Rows: {content.shape[0]} | Columns: {content.shape[1]}")
            else:
                st.text_area(f"Text Content: {fn}", content[:2000], height=200)
else:
    st.info("Please upload a file in the sidebar to begin.")

# --- 7. Chat Interface ---
st.divider()
chat_container = st.container()

with chat_container:
    for msg in st.session_state.history:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])
            if "fig" in msg: st.pyplot(msg["fig"])

if prompt := st.chat_input("Ask a question or give a command (e.g., 'What's the sum of Gold?' or 'Plot sales')"):
    st.session_state.history.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    # Context & Variables setup
    exec_locals = {"pd": pd, "plt": plt, "sns": sns}
    context_str = "You have access to these variables in memory:\n"
    for fn, content in st.session_state.files.items():
        var_name = re.sub(r'[^a-zA-Z0-9]', '_', fn)
        exec_locals[var_name] = content
        if isinstance(content, pd.DataFrame):
            context_str += f"- Variable `{var_name}`: DataFrame from {fn}. Columns: {list(content.columns)}\n"
        else:
            context_str += f"- Variable `{var_name}`: Text from {fn}.\n"

    # Plan -> Execute -> Synthesize Loop
    system_msg = f"""
    {context_str}
    You are a Data Scientist. 
    1. If asked to calculate something, use `print()` for the result. 
    2. If asked to change data, update the specific variable directly.
    3. Always wrap Python code in ```python blocks.
    """

    try:
        # Step 1: AI generates and executes code
        response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "system", "content": system_msg}] + st.session_state.history[-5:]
        )
        ai_raw = response.choices[0].message.content

        captured_output = ""
        fig = None
        if "```python" in ai_raw:
            code = ai_raw.split("```python")[1].split("```")[0].strip()

            # Setup output capture
            old_stdout = sys.stdout
            sys.stdout = io.StringIO()

            try:
                plt.close('all')
                exec(code, {}, exec_locals)
                captured_output = sys.stdout.getvalue()

                # Update the actual files in session state
                for var_name, obj in exec_locals.items():
                    for fn in st.session_state.files.keys():
                        if re.sub(r'[^a-zA-Z0-9]', '_', fn) == var_name:
                            st.session_state.files[fn] = obj

                if plt.get_fignums(): fig = plt.gcf()
            except Exception as e:
                captured_output = f"Error during execution: {e}"
            finally:
                sys.stdout = old_stdout

        # Step 2: Final Synthesis (Conversational Response)
        synthesis_msg = f"""
        User asked: {prompt}
        AI Code was: {ai_raw}
        Code Output was: {captured_output}

        Explain the result to the user clearly. Mention any data changes made.
        """

        final_response = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "system", "content": synthesis_msg}]
        )
        final_text = final_response.choices[0].message.content

        with st.chat_message("assistant"):
            st.markdown(final_text)
            if fig: st.pyplot(fig)

        # Save to history
        entry = {"role": "assistant", "content": final_text}
        if fig: entry["fig"] = fig
        st.session_state.history.append(entry)

        # Force a rerun to update the Previews at the top!
        st.rerun()

    except Exception as e:
        st.error(f"System Error: {e}")